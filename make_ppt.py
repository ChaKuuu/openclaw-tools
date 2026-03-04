# -*- coding: utf-8 -*-
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Get Bilibili data
headers = {'User-Agent': 'Mozilla/5.0'}
r = requests.get('https://api.bilibili.com/x/web-interface/popular?ps=10&pn=1', headers=headers, timeout=10)
data = r.json()

# Video data with analysis
videos = []

if data.get('code') == 0:
    lst = data['data']['list']
    for i, v in enumerate(lst, 1):
        title = v.get('title', 'N/A')
        desc = v.get('desc', '')
        stat = v.get('stat', {})
        danmaku = stat.get('danmaku', 0)
        
        # Generate analysis based on title/content
        if '决战' in title or '夺冠' in title:
            technique = "游戏录屏 + 精彩集锦快切"
            reason = "竞技PK激发观众紧张感，参与感强"
        elif '绝区零' in title or '角色展示' in title:
            technique = "高精度CG展示 + 动态运镜"
            reason = "游戏IP热度高，美术精良吸睛"
        elif '送给' in title or '老弟' in title:
            technique = "情感叙事 + 第一人称视角"
            reason = "情感共鸣，打动人心"
        elif '牧神记' in title:
            technique = "影视切片 + 悬念剪辑"
            reason = "IP续作自带流量，剧情吸引"
        elif '爷爷' in title:
            technique = "生活记录 + 反差萌"
            reason = "代际反差引发好奇，温馨有趣"
        elif 'MrBeast' in title:
            technique = "真人出镜 + 慷慨行为"
            reason = "外国网红正能量，内容稀缺"
        elif '嫦娥' in title or '机甲' in title:
            technique = "动画特效 + 创意混剪"
            reason = "二次元创意，中国风创新"
        elif '4399' in title:
            technique = "品牌叙事 + 情怀杀"
            reason = "集体回忆杀，引发共鸣"
        elif 'MC' in title or '卡牌' in title:
            technique = "游戏玩法展示"
            reason = "经典游戏创新玩法"
        elif '烤鱼' in title:
            technique = "美食探店 + 悬念叙事"
            reason = "地域美食引发好奇"
        else:
            technique = "创意内容"
            reason = "内容有趣"
            
        videos.append({
            'rank': i,
            'title': title,
            'desc': desc[:50] if desc else '',
            'danmaku': danmaku,
            'technique': technique,
            'reason': reason
        })

# Create PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = '2026年3月1日 B站热门视频TOP10'
p.font.size = Pt(48)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
tf2 = sub.text_frame
p2 = tf2.paragraphs[0]
p2.text = '拍摄手法技巧与爆火原因分析'
p2.font.size = Pt(28)
p2.alignment = PP_ALIGN.CENTER

# Content slides
for v in videos:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Rank number
    rank_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1.5), Inches(1.2))
    tf = rank_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"#{v['rank']}"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 69, 69)  # Bilibili red
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = v['title']
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Stats
    stat_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11), Inches(0.5))
    tf = stat_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"弹幕数: {v['danmaku']:,}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(128, 128, 128)
    
    # Technique
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(6), Inches(2))
    tf = tech_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📹 拍摄手法技巧"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    p2 = tf.add_paragraph()
    p2.text = v['technique']
    p2.font.size = Pt(20)
    p2.space_before = Pt(12)
    
    # Reason
    reason_box = slide.shapes.add_textbox(Inches(7), Inches(3), Inches(6), Inches(2))
    tf = reason_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔥 爆火原因"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 69, 69)
    
    p2 = tf.add_paragraph()
    p2.text = v['reason']
    p2.font.size = Pt(20)
    p2.space_before = Pt(12)

# Save
desktop = os.path.join(os.path.expanduser('~'), 'Desktop', 'B站热门视频分析.pptx')
prs.save(desktop)
print(f"已保存到: {desktop}")
